"""
BLOKS PPTX Generator
Usage: python generate.py --input <result.json> --output <file.pptx> [--style consulting]
"""
import json, argparse, sys, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Design System ──────────────────────────────────────────────────────────────

STYLES = {
    "consulting": {
        "bg_dark":    RGBColor(0x1A, 0x1A, 0x2E),   # 딥 네이비
        "bg_light":   RGBColor(0xF5, 0xF5, 0xF7),   # 연한 회색
        "accent":     RGBColor(0xC9, 0xA7, 0x4A),   # 골드
        "accent2":    RGBColor(0x4A, 0x90, 0xD9),   # 블루
        "white":      RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark":  RGBColor(0x1A, 0x1A, 0x2E),
        "text_muted": RGBColor(0x66, 0x66, 0x77),
        "border":     RGBColor(0xDD, 0xDD, 0xEE),
        "font_title": "Pretendard",
        "font_body":  "Pretendard",
    },
    "bloks": {
        "bg_dark":    RGBColor(0x0F, 0x0F, 0x1A),   # BLOKS 딥 다크
        "bg_light":   RGBColor(0xF8, 0xF8, 0xFF),
        "accent":     RGBColor(0x7C, 0x3A, 0xED),   # 네온 퍼플
        "accent2":    RGBColor(0x06, 0xB6, 0xD4),   # 시안
        "white":      RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark":  RGBColor(0x0F, 0x0F, 0x1A),
        "text_muted": RGBColor(0x77, 0x77, 0x99),
        "border":     RGBColor(0xDD, 0xDD, 0xEE),
        "font_title": "Pretendard",
        "font_body":  "Pretendard",
    },
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────────────

def rgb(color): return color

def add_textbox(slide, left, top, width, height, text, size=16, bold=False,
                color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_line(slide, left, top, width, color, thickness=Pt(1)):
    from pptx.util import Emu
    connector = slide.shapes.add_connector(1, left, top, left + width, top)
    connector.line.color.rgb = color
    connector.line.width = thickness
    return connector

def strip_markdown(text):
    """Remove markdown formatting for clean slide text."""
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'#+\s+', '', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    return text.strip()

def split_bullets(text, max_bullets=6):
    """Extract bullet points from markdown text."""
    bullets = []
    for line in text.split('\n'):
        line = line.strip()
        if line.startswith(('- ', '• ', '* ')):
            bullets.append(strip_markdown(line[2:]))
        elif line.startswith(tuple(f'{i}.' for i in range(1, 20))):
            bullets.append(strip_markdown(re.sub(r'^\d+\.\s*', '', line)))
        elif line and not line.startswith('#') and len(bullets) < max_bullets:
            clean = strip_markdown(line)
            if len(clean) > 20:
                bullets.append(clean)
    return bullets[:max_bullets]

# ── Slide Builders ─────────────────────────────────────────────────────────────

def slide_cover(prs, s, title, subtitle, meta, style):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # Dark background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=s["bg_dark"])

    # Accent bar left
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill_color=s["accent"])

    # Gold stripe mid-top
    add_rect(slide, Inches(0.36), Inches(2.8), Inches(8), Pt(2), fill_color=s["accent"])

    # Title
    add_textbox(slide, Inches(0.55), Inches(1.5), Inches(11), Inches(1.2),
                title, size=36, bold=True, color=s["white"])

    # Subtitle
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(3.1), Inches(10), Inches(0.8),
                    subtitle, size=18, color=RGBColor(0xCC, 0xCC, 0xDD))

    # Meta (date, author)
    if meta:
        add_textbox(slide, Inches(0.55), Inches(6.5), Inches(10), Inches(0.5),
                    meta, size=11, color=RGBColor(0x99, 0x99, 0xAA))

    # BLOKS brand mark
    add_textbox(slide, Inches(11.5), Inches(6.8), Inches(1.6), Inches(0.5),
                "BLOKS", size=13, bold=True, color=s["accent"], align=PP_ALIGN.RIGHT)

    return slide


def slide_section_header(prs, s, number, title, desc, style):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=s["bg_dark"])
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill_color=s["accent"])
    add_rect(slide, Inches(0.36), Inches(3.2), Inches(12), Pt(1.5), fill_color=s["accent"])

    if number:
        add_textbox(slide, Inches(0.55), Inches(2.2), Inches(2), Inches(0.8),
                    str(number).zfill(2), size=52, bold=True,
                    color=RGBColor(0x44, 0x44, 0x66))

    add_textbox(slide, Inches(0.55), Inches(3.0), Inches(11), Inches(1.0),
                title, size=32, bold=True, color=s["white"])

    if desc:
        add_textbox(slide, Inches(0.55), Inches(4.2), Inches(10), Inches(0.9),
                    strip_markdown(desc)[:200], size=15,
                    color=RGBColor(0xBB, 0xBB, 0xCC))

    add_textbox(slide, Inches(11.5), Inches(6.8), Inches(1.6), Inches(0.5),
                "BLOKS", size=13, bold=True, color=s["accent"], align=PP_ALIGN.RIGHT)
    return slide


def slide_content(prs, s, title, body_text, slide_num=None, style=None):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Light background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=s["bg_light"])

    # Top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), fill_color=s["bg_dark"])
    add_rect(slide, 0, Inches(0.9), SLIDE_W, Pt(3), fill_color=s["accent"])

    # Title
    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(11.5), Inches(0.75),
                title, size=20, bold=True, color=s["white"])

    # Slide number
    if slide_num:
        add_textbox(slide, Inches(12.4), Inches(0.15), Inches(0.8), Inches(0.5),
                    str(slide_num), size=11, color=RGBColor(0x99, 0x99, 0xAA),
                    align=PP_ALIGN.RIGHT)

    # Body — bullets
    bullets = split_bullets(body_text)

    if bullets:
        y = Inches(1.2)
        for i, bullet in enumerate(bullets):
            # Bullet dot
            add_rect(slide, Inches(0.4), y + Inches(0.12), Inches(0.08), Inches(0.08),
                     fill_color=s["accent"])
            # Text
            add_textbox(slide, Inches(0.65), y, Inches(11.8), Inches(0.55),
                        bullet, size=14, color=s["text_dark"])
            y += Inches(0.65)
    else:
        # Plain paragraph
        clean = strip_markdown(body_text)[:600]
        add_textbox(slide, Inches(0.4), Inches(1.2), Inches(12.4), Inches(5.5),
                    clean, size=13, color=s["text_dark"])

    # Footer
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), fill_color=s["bg_dark"])
    add_textbox(slide, Inches(11), Inches(7.18), Inches(2.2), Inches(0.28),
                "BLOKS", size=9, bold=True, color=s["accent"], align=PP_ALIGN.RIGHT)

    return slide


def slide_two_column(prs, s, title, left_title, left_body, right_title, right_body, slide_num=None):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=s["bg_light"])
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), fill_color=s["bg_dark"])
    add_rect(slide, 0, Inches(0.9), SLIDE_W, Pt(3), fill_color=s["accent"])
    add_textbox(slide, Inches(0.4), Inches(0.1), Inches(11.5), Inches(0.75),
                title, size=20, bold=True, color=s["white"])

    # Left column
    add_rect(slide, Inches(0.3), Inches(1.05), Inches(6.1), Inches(5.9),
             fill_color=s["white"], line_color=s["border"])
    add_rect(slide, Inches(0.3), Inches(1.05), Inches(6.1), Inches(0.5),
             fill_color=s["accent2"])
    add_textbox(slide, Inches(0.45), Inches(1.08), Inches(5.8), Inches(0.44),
                left_title, size=13, bold=True, color=s["white"])
    bullets_l = split_bullets(left_body)
    y = Inches(1.7)
    for b in bullets_l:
        add_rect(slide, Inches(0.5), y + Inches(0.1), Inches(0.07), Inches(0.07),
                 fill_color=s["accent2"])
        add_textbox(slide, Inches(0.7), y, Inches(5.5), Inches(0.5),
                    b, size=12, color=s["text_dark"])
        y += Inches(0.58)

    # Right column
    add_rect(slide, Inches(6.9), Inches(1.05), Inches(6.1), Inches(5.9),
             fill_color=s["white"], line_color=s["border"])
    add_rect(slide, Inches(6.9), Inches(1.05), Inches(6.1), Inches(0.5),
             fill_color=s["accent"])
    add_textbox(slide, Inches(7.05), Inches(1.08), Inches(5.8), Inches(0.44),
                right_title, size=13, bold=True, color=s["white"])
    bullets_r = split_bullets(right_body)
    y = Inches(1.7)
    for b in bullets_r:
        add_rect(slide, Inches(7.1), y + Inches(0.1), Inches(0.07), Inches(0.07),
                 fill_color=s["accent"])
        add_textbox(slide, Inches(7.3), y, Inches(5.5), Inches(0.5),
                    b, size=12, color=s["text_dark"])
        y += Inches(0.58)

    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), fill_color=s["bg_dark"])
    add_textbox(slide, Inches(11), Inches(7.18), Inches(2.2), Inches(0.28),
                "BLOKS", size=9, bold=True, color=s["accent"], align=PP_ALIGN.RIGHT)
    return slide


def slide_metrics(prs, s, title, metrics, slide_num=None):
    """metrics: list of {"label": str, "value": str, "sub": str}"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=s["bg_dark"])
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill_color=s["accent"])
    add_textbox(slide, Inches(0.4), Inches(0.3), Inches(12), Inches(0.7),
                title, size=24, bold=True, color=s["white"])
    add_rect(slide, Inches(0.4), Inches(1.05), Inches(12.5), Pt(1.5), fill_color=s["accent"])

    n = len(metrics)
    card_w = Inches(12.2 / max(n, 1))
    for i, m in enumerate(metrics):
        x = Inches(0.4) + i * card_w + Inches(0.1)
        add_rect(slide, x, Inches(1.3), card_w - Inches(0.2), Inches(2.5),
                 fill_color=RGBColor(0x2A, 0x2A, 0x44),
                 line_color=s["accent"])
        add_textbox(slide, x + Inches(0.15), Inches(1.5), card_w - Inches(0.3), Inches(0.4),
                    m.get("label", ""), size=10, color=RGBColor(0xAA, 0xAA, 0xCC))
        add_textbox(slide, x + Inches(0.1), Inches(1.95), card_w - Inches(0.2), Inches(0.9),
                    m.get("value", ""), size=28, bold=True, color=s["accent"])
        if m.get("sub"):
            add_textbox(slide, x + Inches(0.15), Inches(2.9), card_w - Inches(0.3), Inches(0.35),
                        m.get("sub", ""), size=10, color=RGBColor(0x88, 0x88, 0xAA))

    add_textbox(slide, Inches(11.5), Inches(6.8), Inches(1.6), Inches(0.5),
                "BLOKS", size=13, bold=True, color=s["accent"], align=PP_ALIGN.RIGHT)
    return slide


def slide_closing(prs, s, title, message):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=s["bg_dark"])
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill_color=s["accent"])
    add_rect(slide, Inches(0.36), Inches(3.8), Inches(12), Pt(2), fill_color=s["accent"])

    add_textbox(slide, Inches(0.55), Inches(2.4), Inches(11.5), Inches(1.1),
                title, size=34, bold=True, color=s["white"])
    if message:
        add_textbox(slide, Inches(0.55), Inches(4.0), Inches(10), Inches(1.5),
                    strip_markdown(message)[:300], size=15,
                    color=RGBColor(0xCC, 0xCC, 0xDD))

    add_textbox(slide, Inches(11.5), Inches(6.8), Inches(1.6), Inches(0.5),
                "BLOKS", size=13, bold=True, color=s["accent"], align=PP_ALIGN.RIGHT)
    return slide

# ── Content Parser ─────────────────────────────────────────────────────────────

def parse_sections(text):
    """Split AI output text into sections by headings."""
    sections = []
    current_title = ""
    current_body = []

    for line in text.split('\n'):
        if re.match(r'^#{1,3}\s+', line):
            if current_title or current_body:
                sections.append({
                    "title": current_title,
                    "body": '\n'.join(current_body).strip()
                })
            current_title = re.sub(r'^#{1,3}\s+', '', line).strip()
            current_body = []
        else:
            current_body.append(line)

    if current_title or current_body:
        sections.append({
            "title": current_title,
            "body": '\n'.join(current_body).strip()
        })

    return [s for s in sections if s["title"] or s["body"]]


def extract_metrics_from_text(text):
    """Find metric-like patterns: numbers with % or units."""
    patterns = [
        r'(\d+[\.,]?\d*)\s*([%억만원달러배주초분시간주월년]+)',
        r'(TAM|SAM|SOM|ARR|MAU|NPS|ROI|KPI)\s*[:\-]?\s*(\$?[\d,]+[억만]?[원달러]?)',
    ]
    metrics = []
    for pat in patterns:
        for m in re.finditer(pat, text):
            if len(metrics) >= 5:
                break
            label = m.group(0)[:30]
            val = m.group(1) if len(m.groups()) > 0 else m.group(0)
            metrics.append({"label": label, "value": val})
    return metrics[:5]

# ── Main Builder ───────────────────────────────────────────────────────────────

def build_pptx(result: dict, output_path: str, style_name: str = "consulting"):
    s = STYLES.get(style_name, STYLES["consulting"])

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    project_title = result.get("projectTitle", "BLOKS Report")
    started_at    = result.get("startedAt", "")[:10]
    tasks         = result.get("tasks", [])
    characters    = result.get("characters", [])
    char_map      = {c["id"]: c["name"] for c in characters}
    scenario      = result.get("scenario", "")

    # 1. Cover
    char_names = " · ".join(c["name"] for c in characters[:6])
    slide_cover(prs, s, project_title,
                f"AI 캐릭터 {len(characters)}명 협업 결과 보고서",
                f"{started_at}  |  {char_names}",
                style_name)

    # 2. Summary metrics
    done_count = sum(1 for t in tasks if t["state"] in ("Done", "Approved", "InReview"))
    total_cost = 0
    total_tokens = 0
    for t in tasks:
        ao = t.get("aiOutput")
        if isinstance(ao, str):
            try: ao = json.loads(ao)
            except: ao = None
        if isinstance(ao, dict):
            total_cost   += ao.get("cost_usd", ao.get("costUsd", 0)) or 0
            total_tokens += ao.get("tokens_used", ao.get("tokensUsed", 0)) or 0

    slide_metrics(prs, s, "실행 요약", [
        {"label": "완료 태스크", "value": f"{done_count}/{len(tasks)}", "sub": "100%" if done_count == len(tasks) else ""},
        {"label": "참여 캐릭터", "value": f"{len(characters)}명",      "sub": "AI 전문가"},
        {"label": "총 토큰",    "value": f"{total_tokens:,}",          "sub": "처리량"},
        {"label": "총 비용",    "value": f"${total_cost:.4f}",         "sub": "USD"},
    ])

    # 3. One slide per task
    for idx, task in enumerate(tasks):
        ao = task.get("aiOutput")
        if isinstance(ao, str):
            try: ao = json.loads(ao)
            except: ao = None

        text = ""
        if isinstance(ao, dict):
            text = ao.get("text", ao.get("output_text", ao.get("outputText", ""))) or ""

        title = task.get("title", f"태스크 {idx+1}")
        char_name = char_map.get(task.get("assigneeCharacterId", ""), "AI")

        # Section header
        slide_section_header(prs, s, idx + 1, title, f"담당: {char_name}", style_name)

        if not text:
            continue

        sections = parse_sections(text)

        if len(sections) >= 2 and sections[0]["body"] and sections[1]["body"]:
            # Two-column if first two sections are substantial
            slide_two_column(
                prs, s, title,
                sections[0]["title"] or "주요 내용",
                sections[0]["body"],
                sections[1]["title"] or "세부 사항",
                sections[1]["body"],
                slide_num=idx + 1,
            )
            remaining = sections[2:]
        else:
            remaining = sections

        # Additional content slides (up to 3 per task)
        for sec in remaining[:3]:
            if sec["body"] and len(sec["body"]) > 30:
                slide_content(prs, s,
                              f"{title}  —  {sec['title']}" if sec["title"] else title,
                              sec["body"],
                              slide_num=idx + 1)

    # 4. Closing
    slide_closing(prs, s,
                  "감사합니다",
                  f"본 보고서는 BLOKS AI 캐릭터 {len(characters)}명이 협업하여 작성한 결과물입니다.\n"
                  f"총 {len(tasks)}개 태스크  ·  {total_tokens:,} 토큰  ·  ${total_cost:.4f} USD")

    prs.save(output_path)
    slide_count = len(prs.slides)
    sys.stdout.buffer.write(f"PPTX saved -> {output_path}  ({slide_count} slides)\n".encode("utf-8"))


# ── CLI ────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--input",  required=True, help="result JSON path")
    parser.add_argument("--output", required=True, help="output .pptx path")
    parser.add_argument("--style",  default="consulting", choices=list(STYLES.keys()))
    args = parser.parse_args()

    data = json.loads(Path(args.input).read_text(encoding="utf-8"))
    build_pptx(data, args.output, args.style)
